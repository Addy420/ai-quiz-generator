import os
import google.generativeai as genai
from flask import Flask, request, render_template, redirect, url_for
from PyPDF2 import PdfReader
from pptx import Presentation
from dotenv import load_dotenv

load_dotenv()

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'

# Configure Gemini
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
model = genai.GenerativeModel('gemini-pro')

def extract_text(file_path, filename):
    if filename.endswith('.pdf'):
        text = ""
        with open(file_path, 'rb') as file:
            reader = PdfReader(file)
            for page in reader.pages:
                text += page.extract_text()
        return text
    elif filename.endswith('.pptx'):
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    else:
        return None

def generate_quiz(text):
    prompt = f"""
    Generate 5 multiple-choice questions (MCQs) and 2 short-answer questions from this text:
    {text}

    Format strictly as:
    MCQs:
    1. [Question]?
    a) Option1
    b) Option2
    c) Option3
    d) Option4
    Correct: [Letter]

    Short Answer:
    1. [Question]?
    """

    response = model.generate_content(prompt)
    return response.text.split('\n')

@app.route('/', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        if 'file' not in request.files:
            return redirect(request.url)
        
        file = request.files['file']
        if file.filename == '':
            return redirect(request.url)

        if file:
            filename = file.filename
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)
            
            text = extract_text(file_path, filename)
            if not text:
                return "Unsupported file format"
            
            quiz = generate_quiz(text)
            return render_template('dashboard.html', quiz=quiz)
    
    return render_template('index.html')

if __name__ == '__main__':
    app.run(debug=True)